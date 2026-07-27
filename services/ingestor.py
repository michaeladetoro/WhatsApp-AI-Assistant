import io
import logging
from pathlib import Path

import pytesseract
from pdf2image import convert_from_bytes
from PIL import Image
from pptx import Presentation
from pypdf import PdfReader

logger = logging.getLogger("RAGBot.Ingestor")

# ---------------------------------------------------------------------------
# Supported extensions
# ---------------------------------------------------------------------------
SUPPORTED_EXTENSIONS = {
    ".txt",
    ".md",
    ".pdf",
    ".docx",
    ".docm",
    ".pptx",
    ".pptm",
    ".xlsx",
    ".xlsm",
    ".xlam",
    ".xls",
    ".png",
    ".jpg",
    ".jpeg",
    ".tiff",
    ".bmp",
    ".webp",
}

# Extensions that need OCR (image-based) — subset of PDF handling
_NEED_OCR_HINT = {".pdf"}


class UnsupportedFileTypeError(ValueError):
    pass


class ScannedPDFError(ValueError):
    """Raised when a PDF has no extractable text and OCR is not available."""

    pass


class Ingestor:
    """Stateless document text extractor — all methods are static."""

    @staticmethod
    async def extract(file_bytes: bytes, filename: str) -> str:
        """
        Extract plain text from a file.

        Args:
            file_bytes: Raw file bytes.
            filename:   Original filename (used to detect format by extension).

        Returns:
            Extracted plain text string.

        Raises:
            UnsupportedFileTypeError: Extension not in SUPPORTED_EXTENSIONS.
            ScannedPDFError:          PDF has no text layer and OCR not installed.
            ValueError:               Extraction succeeded but found no content.
        """
        ext = Path(filename).suffix.lower()

        if ext in {".txt", ".md"}:
            return Ingestor._plain_text(file_bytes)

        if ext == ".pdf":
            return Ingestor._pdf(file_bytes, filename)

        if ext in {".docx", ".docm"}:
            return Ingestor._docx(file_bytes)

        if ext in {".pptx", ".pptm"}:
            return Ingestor._pptx(file_bytes)

        if ext in {".xlsx", ".xlsm", ".xlam"}:
            return Ingestor._xlsx(file_bytes)

        if ext == ".xls":
            return Ingestor._xls(file_bytes)

        if ext in {".png", ".jpg", ".jpeg", ".tiff", ".bmp", ".webp"}:
            return Ingestor._image(file_bytes, filename)

        if ext in {".doc", ".ppt"}:
            raise UnsupportedFileTypeError(
                f"Legacy format '{ext}' cannot be parsed directly. "
                "Please re-save as .docx or .pptx and re-upload."
            )

        raise UnsupportedFileTypeError(
            f"File type '{ext}' is not supported.\n"
            f"Supported: {', '.join(sorted(SUPPORTED_EXTENSIONS))}\n"
            "For Google Docs/Sheets/Slides: File → Download → .docx / .xlsx / .pptx"
        )

    # ------------------------------------------------------------------
    # Plain text
    # ------------------------------------------------------------------
    @staticmethod
    def _plain_text(file_bytes: bytes) -> str:
        for enc in ("utf-8", "utf-8-sig", "latin-1", "cp1252"):
            try:
                return file_bytes.decode(enc).strip()
            except UnicodeDecodeError:
                continue
        raise ValueError("Could not decode text file — unknown encoding.")

    # ------------------------------------------------------------------
    # PDF  (pypdf primary → OCR fallback)
    # ------------------------------------------------------------------
    @staticmethod
    def _pdf(file_bytes: bytes, filename: str) -> str:
        """
        Two-stage PDF extraction:
          Stage 1 — pypdf   : instant, zero deps, works for text-based PDFs.
          Stage 2 — OCR     : slow, requires Tesseract + Poppler.
                              Only triggered when Stage 1 yields no text.
        """
        text = Ingestor._pdf_pypdf(file_bytes)

        if text.strip():
            return text

        # No text layer found — try OCR
        logger.info(f"pypdf found no text in '{filename}', attempting OCR fallback...")
        return Ingestor._pdf_ocr(file_bytes, filename)

    @staticmethod
    def _pdf_pypdf(file_bytes: bytes) -> str:

        reader = PdfReader(io.BytesIO(file_bytes))
        pages = [page.extract_text() or "" for page in reader.pages]
        return "\n\n".join(p.strip() for p in pages if p.strip())

    @staticmethod
    def _pdf_ocr(file_bytes: bytes, filename: str) -> str:
        """
        OCR fallback using pdf2image + pytesseract.
        Only invoked when pypdf finds no text (scanned / image-only PDF).
        """

        try:
            pages = convert_from_bytes(file_bytes, dpi=150)
            texts = []
            for i, page_img in enumerate(pages, start=1):
                page_text = pytesseract.image_to_string(page_img).strip()
                if page_text:
                    texts.append(f"[Page {i}]\n{page_text}")
            if not texts:
                raise ValueError(f"OCR found no readable text in '{filename}'.")
            logger.info(f"OCR extracted text from {len(texts)} page(s) of '{filename}'")
            return "\n\n".join(texts)
        except ScannedPDFError:
            raise
        except Exception as e:
            raise ValueError(f"OCR extraction failed for '{filename}': {e}")

    # ------------------------------------------------------------------
    # Word (.docx / .docm)
    # ------------------------------------------------------------------
    @staticmethod
    def _docx(file_bytes: bytes) -> str:
        try:
            from docx import Document
        except ImportError:
            raise ImportError("python-docx is required. Run: pip install python-docx")

        doc = Document(io.BytesIO(file_bytes))
        paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]

        # Also extract text from tables
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(
                    cell.text.strip() for cell in row.cells if cell.text.strip()
                )
                if row_text:
                    paragraphs.append(row_text)

        if not paragraphs:
            raise ValueError("No text found in the Word document.")
        return "\n\n".join(paragraphs)

    # ------------------------------------------------------------------
    # PowerPoint (.pptx / .pptm)
    # ------------------------------------------------------------------
    @staticmethod
    def _pptx(file_bytes: bytes) -> str:

        prs = Presentation(io.BytesIO(file_bytes))
        slides = []
        for i, slide in enumerate(prs.slides, start=1):
            texts = []
            for shape in slide.shapes:
                # Text boxes, titles, body
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
                # Tables embedded in slides
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = " | ".join(
                            cell.text.strip() for cell in row.cells if cell.text.strip()
                        )
                        if row_text:
                            texts.append(row_text)
            if texts:
                slides.append(f"[Slide {i}]\n" + "\n".join(texts))

        if not slides:
            raise ValueError("No text found in the PowerPoint presentation.")
        return "\n\n".join(slides)

    # ------------------------------------------------------------------
    # Excel modern (.xlsx / .xlsm / .xlam)
    # ------------------------------------------------------------------
    @staticmethod
    def _xlsx(file_bytes: bytes) -> str:
        try:
            import openpyxl
        except ImportError:
            raise ImportError("openpyxl is required. Run: pip install openpyxl")

        wb = openpyxl.load_workbook(
            io.BytesIO(file_bytes),
            read_only=True,
            data_only=True,  # Return computed cell values, not formulas
        )
        sheets = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = []
            for row in ws.iter_rows(values_only=True):
                # Skip entirely empty rows
                row_vals = [str(v) for v in row if v is not None and str(v).strip()]
                if row_vals:
                    rows.append(" | ".join(row_vals))
            if rows:
                sheets.append(f"[Sheet: {sheet_name}]\n" + "\n".join(rows))
        wb.close()

        if not sheets:
            raise ValueError("No data found in the spreadsheet.")
        return "\n\n".join(sheets)

    # ------------------------------------------------------------------
    # Excel legacy (.xls)
    # ------------------------------------------------------------------
    @staticmethod
    def _xls(file_bytes: bytes) -> str:
        try:
            import xlrd
        except ImportError:
            raise ImportError("xlrd is required for .xls files. Run: pip install xlrd")

        wb = xlrd.open_workbook(file_contents=file_bytes)
        sheets = []
        for sheet in wb.sheets():
            rows = []
            for rx in range(sheet.nrows):
                row_vals = [
                    str(sheet.cell_value(rx, cx)).strip()
                    for cx in range(sheet.ncols)
                    if str(sheet.cell_value(rx, cx)).strip()
                ]
                if row_vals:
                    rows.append(" | ".join(row_vals))
            if rows:
                sheets.append(f"[Sheet: {sheet.name}]\n" + "\n".join(rows))

        if not sheets:
            raise ValueError("No data found in the legacy spreadsheet.")
        return "\n\n".join(sheets)

    # ------------------------------------------------------------------
    # Images (OCR)
    # ------------------------------------------------------------------
    @staticmethod
    def _image(file_bytes: bytes, filename: str) -> str:
        """
        Extract text from images using OCR (pytesseract + Pillow).
        """

        try:
            image = Image.open(io.BytesIO(file_bytes))
            text = pytesseract.image_to_string(image).strip()
            if not text:
                raise ValueError(
                    f"OCR found no readable text in the image '{filename}'."
                )
            logger.info(f"OCR extracted text from '{filename}'")
            return text
        except UnsupportedFileTypeError:
            raise
        except Exception as e:
            raise ValueError(f"Image extraction failed for '{filename}': {e}")
